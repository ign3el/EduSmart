import os
import json
import random
import io
import wave
from google import genai
from google.genai import types
import docx
import pptx

class GeminiService:
    def __init__(self):
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            print("CRITICAL WARNING: GEMINI_API_KEY is missing!")
        
        self.client = genai.Client(api_key=api_key)
        
        # 1. Text: Gemini 2.5 Flash
        self.text_model = "gemini-2.5-flash"
        
        # 2. Image: Imagen 4.0 (Working!)
        self.image_model = "imagen-4.0-generate-001" 
        
        # 3. Audio: Switched to specific dated version for stability
        self.audio_model = "gemini-2.5-flash-native-audio-preview-12-2025"

    def process_file_to_story(self, file_path: str):
        ext = os.path.splitext(file_path)[1].lower()
        content_part = None
        
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                pdf_data = f.read()
            content_part = types.Part.from_bytes(data=pdf_data, mime_type="application/pdf")
        elif ext == ".docx":
            try:
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                content_part = f"Word Doc Content:\n{text}"
            except: return None
        elif ext == ".pptx":
            try:
                prs = pptx.Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                content_part = f"PowerPoint Content:\n" + "\n".join(text)
            except: return None
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                content_part = f.read()

        if not content_part:
            return None

        try:
            prompt = """
            You are a game designer converting this content into an Interactive Visual Novel for kids (Age 6-8).
            Create a linear story where the user learns the topic through a narrative.
            
            Return strictly valid JSON (no markdown) with this structure:
            {
              "title": "Story Title",
              "story_id": "unique_id_placeholder",
              "scenes": [
                {
                  "id": 1,
                  "text": "Story narration for this scene (engaging, max 2 sentences).",
                  "image_description": "Visual description for the artist, 3d pixar style",
                  "character_details": "Blue robot named Beep",
                  "interaction": {
                      "question": "A simple question about what just happened?",
                      "options": ["Option A", "Option B"],
                      "correct_answer": "Option A"
                  }
                }
              ]
            }
            """
            contents = [content_part, prompt]
            response = self.client.models.generate_content(
                model=self.text_model,
                contents=contents,
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )
            return json.loads(response.text)
        except Exception as e:
            print(f"Story Gen Error: {e}")
            return None

    def generate_image(self, prompt: str, seed: int = None):
        """Generates Images using Imagen 4.0."""
        try:
            if seed is None:
                seed = random.randint(0, 2**32 - 1)
            
            response = self.client.models.generate_images(
                model=self.image_model,
                prompt=f"kids educational illustration, 3d pixar style, vibrant: {prompt}",
                config=types.GenerateImagesConfig(
                    number_of_images=1,
                )
            )
            return response.generated_images[0].image.image_bytes
        except Exception as e:
            print(f"Image Gen Error: {e}")
            return None

    def generate_voiceover(self, text: str):
        """Generates Audio (With Silent Fallback if API fails)."""
        try:
            response = self.client.models.generate_content(
                model=self.audio_model,
                contents=f"Narrate this for a child in a cheerful, energetic voice: {text}",
                config=types.GenerateContentConfig(
                    response_modalities=["AUDIO"]
                )
            )
            for part in response.candidates[0].content.parts:
                if part.inline_data:
                    return part.inline_data.data
            
            # If no audio data found, trigger fallback
            raise Exception("No inline audio data returned")

        except Exception as e:
            print(f"Audio Gen Error: {e}")
            # FALLBACK: Generate 1 second of silence so the app doesn't break
            return self._generate_silent_wav()

    def _generate_silent_wav(self):
        """Helper to create a dummy silent WAV file."""
        buffer = io.BytesIO()
        with wave.open(buffer, 'wb') as wav:
            wav.setnchannels(1) # Mono
            wav.setsampwidth(2) # 16-bit
            wav.setframerate(44100)
            wav.writeframes(b'\x00' * 44100) # 1 second of silence
        return buffer.getvalue()